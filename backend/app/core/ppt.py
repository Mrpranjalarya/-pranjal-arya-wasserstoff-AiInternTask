# 2. backend/app/core/ppt.py
import os
from pptx import Presentation

def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract text from a .pptx file, tagging slides.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    prs = Presentation(file_path)
    output = ""
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    texts.append(t)
        if texts:
            slide_text = "\n".join(texts)
            output += f"\n--- Slide {idx} ---\n{slide_text}\n"
    return output.strip()
